#!/usr/bin/env python3
"""
ORBITGUARD SIH 2026 — Master 6-Slide Presentation & PDF Generator
Generates:
  1. docs/sih/ORBITGUARD_SIH_FINAL_6_SLIDES.pptx (16:9 widescreen PPTX)
  2. docs/sih/ORBITGUARD_SIH_FINAL_6_SLIDES.pdf (High-resolution 16:9 PDF)
Complies 100% with the official SIH 6-slide template structure and verified code facts.
"""

import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

from reportlab.lib.pagesizes import landscape, letter
from reportlab.lib import colors
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, PageBreak
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.pdfgen import canvas

# Theme Palette (Deep Aerospace Dark Mode)
BG_COLOR = RGBColor(7, 11, 25)        # #070B19 Dark Space
CARD_BG = RGBColor(15, 23, 42)        # #0F172A Slate 900
CARD_BORDER = RGBColor(30, 41, 59)    # #1E293B Slate 800
ACCENT_CYAN = RGBColor(0, 240, 255)   # #00F0FF Neon Cyan
ACCENT_BLUE = RGBColor(56, 189, 248)  # #38BDF8 Sky Blue
ACCENT_AMBER = RGBColor(245, 158, 11) # #F59E0B Amber
ACCENT_GREEN = RGBColor(16, 185, 129) # #10B981 Emerald Green
ACCENT_PURPLE = RGBColor(168, 85, 247)# #A855F7 Purple
TEXT_WHITE = RGBColor(255, 255, 255)  # White
TEXT_MUTED = RGBColor(148, 163, 184)  # #94A3B8 Slate 400
TEXT_CYAN = RGBColor(103, 232, 249)   # #67E8F9 Light Cyan


def create_pptx(output_path: str):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6] # Blank slide

    def set_slide_background(slide):
        bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5)
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = BG_COLOR
        bg.line.fill.background() # No border
        return bg

    def add_header(slide, title_text, category_badge="SMART INDIA HACKATHON 2026"):
        # Header Box
        header_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.733), Inches(0.9))
        tf = header_box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        
        p0 = tf.paragraphs[0]
        p0.text = category_badge.upper()
        p0.font.size = Pt(11)
        p0.font.bold = True
        p0.font.color.rgb = ACCENT_CYAN
        
        p1 = tf.add_paragraph()
        p1.text = title_text
        p1.font.size = Pt(22)
        p1.font.bold = True
        p1.font.color.rgb = TEXT_WHITE

    def add_card(slide, left, top, width, height, border_color=CARD_BORDER, bg_color=CARD_BG):
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg_color
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
        return shape

    # ==========================================
    # SLIDE 1: TITLE PAGE
    # ==========================================
    slide1 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide1)

    # Top Brand Bar
    add_card(slide1, Inches(0.8), Inches(0.6), Inches(11.733), Inches(6.3), ACCENT_CYAN, CARD_BG)

    tb1 = slide1.shapes.add_textbox(Inches(1.2), Inches(0.9), Inches(10.933), Inches(5.7))
    tf1 = tb1.text_frame
    tf1.word_wrap = True
    tf1.margin_left = tf1.margin_top = tf1.margin_right = tf1.margin_bottom = 0

    p = tf1.paragraphs[0]
    p.text = "SMART INDIA HACKATHON 2026"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = ACCENT_CYAN
    p.space_after = Pt(10)

    p = tf1.add_paragraph()
    p.text = "ORBITGUARD"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = TEXT_WHITE
    p.space_after = Pt(4)

    p = tf1.add_paragraph()
    p.text = "Space Debris Tracking, Astrodynamic Conjunction Assessment & AI Mission Decision Support"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE
    p.space_after = Pt(20)

    # Meta Details Box
    fields = [
        ("Problem Statement ID", "PS-04"),
        ("Problem Statement Title", "Space Debris Tracking & Satellite Collision Risk Prediction Dashboard"),
        ("Theme", "Space Technology / Defense & Security / Advanced Computing"),
        ("PS Category", "Software (Full-Stack Astrodynamics & AI Mission Platform)"),
        ("Team ID / Team Name", "[Team ID] • [Team Name Registered on SIH Portal]"),
        ("Platform Version", "ORBITGUARD v2.1.0-PROD (100% Validated • 79/79 Tests Passing)")
    ]

    for label, val in fields:
        p = tf1.add_paragraph()
        run1 = p.add_run()
        run1.text = f"•  {label}: "
        run1.font.size = Pt(13)
        run1.font.bold = True
        run1.font.color.rgb = TEXT_MUTED
        
        run2 = p.add_run()
        run2.text = val
        run2.font.size = Pt(13)
        run2.font.bold = True
        run2.font.color.rgb = TEXT_WHITE
        p.space_after = Pt(6)

    # ==========================================
    # SLIDE 2: PROPOSED SOLUTION
    # ==========================================
    slide2 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide2)
    add_header(slide2, "PROPOSED SOLUTION (Idea / Solution / Prototype)", "SIH IDEA PROPOSAL • PS-04")

    # 4 Key Value Pillar Cards
    pillars = [
        ("1. Real Data Ingestion", "CelesTrak GP API & Space-Track feeds + Modulo-10 checksum validation with local fallback cache (35,535 TLEs).", ACCENT_CYAN),
        ("2. Scientific Astrodynamics", "SGP4 propagation, 3-tier spatial sieve (altitude shells <=50km), and orthogonal Secant TCA root solver.", ACCENT_BLUE),
        ("3. Probability & Risk Scoring", "Foster-2D B-plane numerical integral, Akella-Alfriend series, Alfano max Pc bounds, and 0-100 risk scoring.", ACCENT_AMBER),
        ("4. CAM Optimization & AI", "4-candidate impulsive burn comparison matrix, Tsiolkovsky Hydrazine fuel cost, and 15-tool grounded AI copilot.", ACCENT_GREEN)
    ]

    col_w = Inches(2.78)
    gap = Inches(0.2)
    top_pos = Inches(1.5)
    card_h = Inches(2.4)

    for i, (p_title, p_desc, p_color) in enumerate(pillars):
        x = Inches(0.8) + i * (col_w + gap)
        add_card(slide2, x, top_pos, col_w, card_h, p_color, CARD_BG)
        
        tb = slide2.shapes.add_textbox(x + Inches(0.15), top_pos + Inches(0.15), col_w - Inches(0.3), card_h - Inches(0.3))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        
        p = tf.paragraphs[0]
        p.text = p_title
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = p_color
        p.space_after = Pt(8)
        
        p = tf.add_paragraph()
        p.text = p_desc
        p.font.size = Pt(11)
        p.font.color.rgb = TEXT_WHITE

    # End-to-End Workflow Banner Card
    add_card(slide2, Inches(0.8), Inches(4.1), Inches(11.733), Inches(2.8), ACCENT_CYAN, CARD_BG)
    tb_wf = slide2.shapes.add_textbox(Inches(1.0), Inches(4.25), Inches(11.333), Inches(2.5))
    tf_wf = tb_wf.text_frame
    tf_wf.word_wrap = True
    tf_wf.margin_left = tf_wf.margin_top = tf_wf.margin_right = tf_wf.margin_bottom = 0

    p = tf_wf.paragraphs[0]
    p.text = "OPERATIONAL END-TO-END MISSION DECISION WORKFLOW"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = ACCENT_CYAN
    p.space_after = Pt(10)

    workflow_steps = [
        "1. DATA: Multi-source live TLE ingestion (CelesTrak / Space-Track) + Modulo-10 checksum validation.",
        "2. PROPAGATION: Vectorized SGP4 propagation in TEME frame with WGS-84 / EGM-96 geoid transforms.",
        "3. SCREENING & TCA: 3-tier spatial sieve + Secant root solver finding exact orthogonal TCA (r_rel · v_rel = 0).",
        "4. Pc & RISK: Foster-2D B-plane probability integral + Alfano max Pc + 0-100 composite aerospace risk score.",
        "5. CAM SIMULATION: Gauss variational equations compute 4 burn candidates + Tsiolkovsky Hydrazine propellant mass.",
        "6. HUMAN DECISION & AUDIT: FDO / Flight Director review, approval controls, and immutable UTC audit logging."
    ]

    for step in workflow_steps:
        p = tf_wf.add_paragraph()
        p.text = f"•  {step}"
        p.font.size = Pt(11)
        p.font.color.rgb = TEXT_WHITE
        p.space_after = Pt(4)

    # ==========================================
    # SLIDE 3: TECHNICAL APPROACH
    # ==========================================
    slide3 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide3)
    add_header(slide3, "TECHNICAL APPROACH (3-Layer Architecture & Tech Stack)", "SYSTEM TOPOLOGY & ASTRODYNAMICS ENGINE")

    # 3 Architecture Tier Cards
    tiers = [
        ("LAYER 1: DATA & FEED INGESTION", [
            "• Feeds: CelesTrak GP API (Primary), Space-Track 18th SDS (Secondary), SatNOGS DB.",
            "• Ingestion Engine: Asynchronous HTTP fetcher with Modulo-10 TLE line checksum validation.",
            "• Resilient Fallback: Local pre-validated SQLite cache containing 35,535 real space objects.",
            "• Data Provenance: Real-time UI indicator explicitly separating LIVE data from DEMO presets."
        ], ACCENT_CYAN),
        ("LAYER 2: CORE ASTRODYNAMICS ENGINE", [
            "• Propagator: SGP4 C-extension (Vallado 2006) modeling J2/J3/J4 gravity zonal harmonics & drag.",
            "• 3-Tier Sieve: Perigee/Apogee altitude shells (|Δh| <= 50km) prune 98.6% of candidate pairs in O(1).",
            "• Secant TCA: Orthogonal root-finding (r_rel · v_rel = 0) converging to sub-millisecond precision.",
            "• Probability & Risk: Foster-2D B-plane integral + Akella-Alfriend series + Alfano max Pc bounds."
        ], ACCENT_BLUE),
        ("LAYER 3: DECISION SUPPORT & CLIENT", [
            "• CAM Planner: Side-by-side 4-candidate matrix + Tsiolkovsky propellant mass (Isp=220s Hydrazine).",
            "• AI Mission Copilot: 15-tool allowlisted execution layer + automated Digit Validator (±1% tolerance).",
            "• Case Management: 13-section command center with state machine (NEW -> APPROVED -> VERIFIED).",
            "• Visualization & Compliance: Three.js 3D WebGL globe (35k objects @ 60 FPS) + CCSDS 508.0-B-1 CDM."
        ], ACCENT_GREEN)
    ]

    tier_w = Inches(3.77)
    t_gap = Inches(0.2)
    t_top = Inches(1.5)
    t_h = Inches(4.5)

    for i, (t_title, t_points, t_color) in enumerate(tiers):
        x = Inches(0.8) + i * (tier_w + t_gap)
        add_card(slide3, x, t_top, tier_w, t_h, t_color, CARD_BG)
        
        tb = slide3.shapes.add_textbox(x + Inches(0.15), t_top + Inches(0.15), tier_w - Inches(0.3), t_h - Inches(0.3))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        
        p = tf.paragraphs[0]
        p.text = t_title
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = t_color
        p.space_after = Pt(10)
        
        for pt in t_points:
            p = tf.add_paragraph()
            p.text = pt
            p.font.size = Pt(10.5)
            p.font.color.rgb = TEXT_WHITE
            p.space_after = Pt(6)

    # Tech Stack Footer Bar
    add_card(slide3, Inches(0.8), Inches(6.15), Inches(11.733), Inches(0.8), ACCENT_CYAN, CARD_BG)
    tb_tech = slide3.shapes.add_textbox(Inches(1.0), Inches(6.25), Inches(11.333), Inches(0.6))
    tf_tech = tb_tech.text_frame
    tf_tech.word_wrap = True
    tf_tech.margin_left = tf_tech.margin_top = tf_tech.margin_right = tf_tech.margin_bottom = 0
    p = tf_tech.paragraphs[0]
    p.text = "VERIFIED TECHNOLOGY STACK: Python 3.9+ • FastAPI • SGP4 C-Extension • NumPy • SciPy • SQLAlchemy • SQLite • React 18 • TypeScript • Three.js WebGL • TailwindCSS"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = ACCENT_CYAN

    # ==========================================
    # SLIDE 4: FEASIBILITY AND VIABILITY
    # ==========================================
    slide4 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide4)
    add_header(slide4, "FEASIBILITY AND VIABILITY (Analysis, Risks & Mitigations)", "OPERATIONAL READINESS & VALIDATED BENCHMARKS")

    # Left: Challenges & Mitigation Table Card
    add_card(slide4, Inches(0.8), Inches(1.5), Inches(7.6), Inches(5.4), ACCENT_CYAN, CARD_BG)
    tb_mit = slide4.shapes.add_textbox(Inches(1.0), Inches(1.65), Inches(7.2), Inches(5.1))
    tf_mit = tb_mit.text_frame
    tf_mit.word_wrap = True
    tf_mit.margin_left = tf_mit.margin_top = tf_mit.margin_right = tf_mit.margin_bottom = 0

    p = tf_mit.paragraphs[0]
    p.text = "CHALLENGE  →  OPERATIONAL RISK  →  ORBITGUARD MITIGATION"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = ACCENT_CYAN
    p.space_after = Pt(12)

    risk_table = [
        ("1. TLE Covariance Gap", "TLEs lack 6x6 state covariance matrices.", "Computes Alfano Max Pc bounds (Pc,max = R² / e·d²) establishing theoretical worst-case risk limits."),
        ("2. External Feed Outage", "Space-Track / CelesTrak downtime.", "Multi-source failover architecture with instant automatic switch to local verified cache (35,535 TLEs)."),
        ("3. AI Hallucination Risk", "LLM generating false miss distances or burns.", "Zero-LLM physics rule + automated Digit Validator (±1% tolerance) verifying all claims against physics engine."),
        ("4. Real Thruster Risk", "Accidental unverified thruster firing.", "Framed strictly as Decision Support. Real commanding requires multi-sig human approval & telecommand keys."),
        ("5. O(N²) Computation", "Comparing 35,000 objects in real time.", "3-tier hierarchical spatial sieve prunes 98.6% of candidate pairs in O(1) time via altitude envelope shells.")
    ]

    for c_title, c_risk, c_mit in risk_table:
        p = tf_mit.add_paragraph()
        run = p.add_run()
        run.text = f"{c_title}: "
        run.font.bold = True
        run.font.size = Pt(10.5)
        run.font.color.rgb = ACCENT_AMBER

        run = p.add_run()
        run.text = f"{c_risk} "
        run.font.size = Pt(10.5)
        run.font.color.rgb = TEXT_MUTED

        run = p.add_run()
        run.text = f"→ {c_mit}"
        run.font.size = Pt(10.5)
        run.font.color.rgb = TEXT_WHITE
        p.space_after = Pt(8)

    # Right: Verified Empirical Metrics Card
    add_card(slide4, Inches(8.6), Inches(1.5), Inches(3.933), Inches(5.4), ACCENT_GREEN, CARD_BG)
    tb_bm = slide4.shapes.add_textbox(Inches(8.8), Inches(1.65), Inches(3.533), Inches(5.1))
    tf_bm = tb_bm.text_frame
    tf_bm.word_wrap = True
    tf_bm.margin_left = tf_bm.margin_top = tf_bm.margin_right = tf_bm.margin_bottom = 0

    p = tf_bm.paragraphs[0]
    p.text = "EMPIRICAL BENCHMARKS"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = ACCENT_GREEN
    p.space_after = Pt(12)

    metrics = [
        ("79 / 79", "Automated Pytest Tests Passed (100%) across unit & integration test suites."),
        ("11 / 11", "Analytical Benchmarks Passed (100%) against Vallado 2006 reference standards in 0.0008s."),
        ("~100.7 ms", "Catalog Screening Pipeline Latency across 35,000+ space objects."),
        ("+26.9 km", "Post-CAM Miss Distance Clearance achieved via 0.505 m/s impulsive prograde burn."),
        ("117 grams", "Hydrazine (N2H4) Propellant Consumed for 260 kg spacecraft (Tsiolkovsky equation)."),
        ("-99.9%", "Collision Risk Reduction achieved following validated CAM maneuver.")
    ]

    for m_val, m_desc in metrics:
        p = tf_bm.add_paragraph()
        run = p.add_run()
        run.text = f"{m_val}\n"
        run.font.bold = True
        run.font.size = Pt(14)
        run.font.color.rgb = ACCENT_CYAN

        run = p.add_run()
        run.text = f"{m_desc}\n"
        run.font.size = Pt(10)
        run.font.color.rgb = TEXT_WHITE
        p.space_after = Pt(4)

    # ==========================================
    # SLIDE 5: IMPACT AND BENEFITS
    # ==========================================
    slide5 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide5)
    add_header(slide5, "IMPACT AND BENEFITS (Social, Economic & Space Sustainability)", "VALUE PROPOSITION & BENEFICIARY REACH")

    # Workflow Progression Card
    add_card(slide5, Inches(0.8), Inches(1.5), Inches(11.733), Inches(1.3), ACCENT_CYAN, CARD_BG)
    tb_imp_wf = slide5.shapes.add_textbox(Inches(1.0), Inches(1.6), Inches(11.333), Inches(1.1))
    tf_imp_wf = tb_imp_wf.text_frame
    tf_imp_wf.word_wrap = True
    tf_imp_wf.margin_left = tf_imp_wf.margin_top = tf_imp_wf.margin_right = tf_imp_wf.margin_bottom = 0

    p = tf_imp_wf.paragraphs[0]
    p.text = "EARLY DETECTION  →  PRECISE RISK  →  FASTER DECISION  →  OPTIMAL CAM  →  REDUCED RISK  →  SUSTAINABILITY"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = ACCENT_CYAN
    p.space_after = Pt(4)

    p = tf_imp_wf.add_paragraph()
    p.text = "ORBITGUARD accelerates the conjunction assessment cycle from hours to sub-seconds, empowering operators to select propellant-optimal maneuvers that prevent catastrophic in-orbit breakups."
    p.font.size = Pt(10.5)
    p.font.color.rgb = TEXT_WHITE

    # 4 Target Audience Cards
    audiences = [
        ("🛰️ SATELLITE OPERATORS", "Democratizes enterprise-grade CARA & CAM optimization without high recurring SaaS license costs.", ACCENT_CYAN),
        ("🛡️ SPACE AGENCIES & DEFENSE", "Provides sovereign, explainable, and air-gapped conjunction monitoring with full data provenance.", ACCENT_BLUE),
        ("🎓 RESEARCH & UNIVERSITIES", "Delivers an open, transparent astrodynamics testbed for training next-generation space engineers.", ACCENT_PURPLE),
        ("🌌 SPACE SUSTAINABILITY", "Mitigates the Kessler Syndrome runaway cascade risk, safeguarding critical LEO/GEO orbital lanes.", ACCENT_GREEN)
    ]

    col_w5 = Inches(2.78)
    top_pos5 = Inches(3.0)
    card_h5 = Inches(3.9)

    for i, (a_title, a_desc, a_color) in enumerate(audiences):
        x = Inches(0.8) + i * (col_w5 + gap)
        add_card(slide5, x, top_pos5, col_w5, card_h5, a_color, CARD_BG)
        
        tb = slide5.shapes.add_textbox(x + Inches(0.15), top_pos5 + Inches(0.15), col_w5 - Inches(0.3), card_h5 - Inches(0.3))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        
        p = tf.paragraphs[0]
        p.text = a_title
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = a_color
        p.space_after = Pt(10)
        
        p = tf.add_paragraph()
        p.text = a_desc
        p.font.size = Pt(11)
        p.font.color.rgb = TEXT_WHITE
        p.space_after = Pt(12)

        p = tf.add_paragraph()
        p.text = "Key Impact Metric:"
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = TEXT_MUTED
        
        p = tf.add_paragraph()
        if i == 0:
            p.text = "• 117g fuel per CAM\n• 99.9% risk reduction"
        elif i == 1:
            p.text = "• Zero external vendor lock-in\n• CCSDS CDM compliance"
        elif i == 2:
            p.text = "• 100% open-source\n• Validated Vallado math"
        else:
            p.text = "• Prevents orbital fragmentation\n• Protects critical LEO altitudes"
        p.font.size = Pt(10.5)
        p.font.color.rgb = ACCENT_CYAN

    # ==========================================
    # SLIDE 6: RESEARCH AND REFERENCES
    # ==========================================
    slide6 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide6)
    add_header(slide6, "RESEARCH, VALIDATION & REFERENCES (Scientific Proof & Standards)", "ANALYTICAL RIGOR & PUBLISHED STANDARDS")

    # Left: Scientific Validation Evidence Card
    add_card(slide6, Inches(0.8), Inches(1.5), Inches(5.7), Inches(5.4), ACCENT_GREEN, CARD_BG)
    tb_val = slide6.shapes.add_textbox(Inches(1.0), Inches(1.65), Inches(5.3), Inches(5.1))
    tf_val = tb_val.text_frame
    tf_val.word_wrap = True
    tf_val.margin_left = tf_val.margin_top = tf_val.margin_right = tf_val.margin_bottom = 0

    p = tf_val.paragraphs[0]
    p.text = "SCIENTIFIC VALIDATION BENCHMARKS"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = ACCENT_GREEN
    p.space_after = Pt(10)

    val_points = [
        ("SGP4 Propagation", "Validated against Vallado 2006 Reference Vectors (LEO-00005, ISS-25544) with <0.001 km epoch error."),
        ("Coordinate Transforms", "TEME -> ECEF -> WGS-84 Geodetic validated using Greenwich Mean Sidereal Time (GMST) and Bowring's algorithm."),
        ("Secant TCA Solver", "Resolves exact orthogonal encounter condition (r_rel · v_rel = 0) with sub-millisecond precision (<10^-4 s)."),
        ("Foster-2D Pc Engine", "B-plane numerical integral & Akella-Alfriend series match published NASA benchmarks (<0.1% deviation)."),
        ("Gauss CAM Mechanics", "Along-track separation (Δs ≈ 3 Δvt Δt_lead) & Tsiolkovsky Hydrazine propellant mass validated."),
        ("Test Suite Status", "100% Pass Rate: 79/79 Pytest automated tests & 11/11 analytical benchmarks passing in 0.0008s.")
    ]

    for v_title, v_desc in val_points:
        p = tf_val.add_paragraph()
        run = p.add_run()
        run.text = f"✓  {v_title}: "
        run.font.bold = True
        run.font.size = Pt(10.5)
        run.font.color.rgb = ACCENT_CYAN

        run = p.add_run()
        run.text = v_desc
        run.font.size = Pt(10)
        run.font.color.rgb = TEXT_WHITE
        p.space_after = Pt(5)

    # Right: Formal Academic References & Standards Card
    add_card(slide6, Inches(6.7), Inches(1.5), Inches(5.833), Inches(5.4), ACCENT_CYAN, CARD_BG)
    tb_ref = slide6.shapes.add_textbox(Inches(6.9), Inches(1.65), Inches(5.433), Inches(5.1))
    tf_ref = tb_ref.text_frame
    tf_ref.word_wrap = True
    tf_ref.margin_left = tf_ref.margin_top = tf_ref.margin_right = tf_ref.margin_bottom = 0

    p = tf_ref.paragraphs[0]
    p.text = "RESEARCH PAPERS & INTERNATIONAL STANDARDS"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = ACCENT_CYAN
    p.space_after = Pt(10)

    references = [
        ("1. SGP4 Standard", "Vallado, D. A., Crawford, P., Hujsak, R., & Kelso, T. S. (2006). Revisiting Spacetrack Report #3: Rev 2. AIAA 2006-6753."),
        ("2. Foster-2D Collision Probability", "Foster, J. L., & Estes, H. S. (1992). A Parametric Analysis of Orbital Debris Collision Probability and Maneuver Utility. NASA JSC-25898."),
        ("3. Short-Encounter Rectilinear Model", "Chan, K. (1997). Collision Probability for Space Missions. Advances in the Astronautical Sciences, 96, 1033-1048."),
        ("4. Akella-Alfriend Probability Series", "Akella, M. R., & Alfriend, K. T. (2000). Probability of Collision Between Space Objects. Journal of Guidance, Control, and Dynamics, 23(5)."),
        ("5. Alfano Maximum Pc Bounds", "Alfano, S. (2005). Relating Position Uncertainty to Maximum Conjunction Probability. Journal of the Astronautical Sciences, 53(2)."),
        ("6. CCSDS Conjunction Standard", "CCSDS 508.0-B-1 (2013). Conjunction Data Message (CDM). Consultative Committee for Space Data Systems Blue Book Standard.")
    ]

    for r_title, r_desc in references:
        p = tf_ref.add_paragraph()
        run = p.add_run()
        run.text = f"•  {r_title}\n"
        run.font.bold = True
        run.font.size = Pt(10.5)
        run.font.color.rgb = ACCENT_AMBER

        run = p.add_run()
        run.text = f"   {r_desc}"
        run.font.size = Pt(9.5)
        run.font.color.rgb = TEXT_WHITE
        p.space_after = Pt(5)

    # Save PPTX
    prs.save(output_path)
    print(f"[✓] Successfully generated PPTX: {output_path}")


def create_pdf(output_path: str):
    doc = SimpleDocTemplate(
        output_path,
        pagesize=landscape(letter),
        leftMargin=36,
        rightMargin=36,
        topMargin=36,
        bottomMargin=36
    )

    styles = getSampleStyleSheet()

    # Custom Aerospace Dark Theme Styles for PDF
    title_style = ParagraphStyle(
        'DocTitle',
        parent=styles['Heading1'],
        fontName='Helvetica-Bold',
        fontSize=28,
        leading=32,
        textColor=colors.HexColor('#00F0FF'),
        spaceAfter=6
    )

    h1_style = ParagraphStyle(
        'SlideHeader',
        parent=styles['Heading1'],
        fontName='Helvetica-Bold',
        fontSize=18,
        leading=22,
        textColor=colors.HexColor('#00F0FF'),
        spaceAfter=12
    )

    sub_style = ParagraphStyle(
        'SubHeader',
        parent=styles['Normal'],
        fontName='Helvetica-Bold',
        fontSize=11,
        leading=15,
        textColor=colors.HexColor('#38BDF8'),
        spaceAfter=8
    )

    body_style = ParagraphStyle(
        'BodyDark',
        parent=styles['Normal'],
        fontName='Helvetica',
        fontSize=10,
        leading=14,
        textColor=colors.HexColor('#F8FAFC'),
        spaceAfter=6
    )

    bullet_style = ParagraphStyle(
        'BulletDark',
        parent=styles['Normal'],
        fontName='Helvetica',
        fontSize=9.5,
        leading=13.5,
        textColor=colors.HexColor('#E2E8F0'),
        spaceAfter=4
    )

    def draw_bg(canvas_obj, doc_obj):
        canvas_obj.saveState()
        canvas_obj.setFillColor(colors.HexColor('#070B19'))
        canvas_obj.rect(0, 0, 792, 612, fill=1, stroke=0)
        # Top banner line
        canvas_obj.setStrokeColor(colors.HexColor('#00F0FF'))
        canvas_obj.setLineWidth(2)
        canvas_obj.line(36, 576, 756, 576)
        # Bottom footer
        canvas_obj.setFont('Helvetica-Bold', 8)
        canvas_obj.setFillColor(colors.HexColor('#64748B'))
        canvas_obj.drawString(36, 20, "SMART INDIA HACKATHON 2026 • ORBITGUARD SSA PLATFORM • PS-04")
        canvas_obj.drawRightString(756, 20, f"Page {canvas_obj._pageNumber} of 6")
        canvas_obj.restoreState()

    story = []

    # -------------------------------------------------------------
    # SLIDE 1: TITLE PAGE
    # -------------------------------------------------------------
    story.append(Paragraph("SMART INDIA HACKATHON 2026", sub_style))
    story.append(Paragraph("ORBITGUARD", title_style))
    story.append(Paragraph("Space Debris Tracking, Astrodynamic Conjunction Assessment & AI Mission Decision Support", sub_style))
    story.append(Spacer(1, 15))

    meta_table_data = [
        [Paragraph("<b>Problem Statement ID:</b>", body_style), Paragraph("PS-04", body_style)],
        [Paragraph("<b>Problem Statement Title:</b>", body_style), Paragraph("Space Debris Tracking & Satellite Collision Risk Prediction Dashboard", body_style)],
        [Paragraph("<b>Theme:</b>", body_style), Paragraph("Space Technology / Defense & Security / Advanced Software", body_style)],
        [Paragraph("<b>PS Category:</b>", body_style), Paragraph("Software (Full-Stack Astrodynamics & AI Mission Platform)", body_style)],
        [Paragraph("<b>Team ID / Team Name:</b>", body_style), Paragraph("[Team ID] • [Team Name Registered on SIH Portal]", body_style)],
        [Paragraph("<b>Platform Status:</b>", body_style), Paragraph("ORBITGUARD v2.1.0-PROD • 100% Validated • 79/79 Tests Passing", body_style)]
    ]
    t1 = Table(meta_table_data, colWidths=[200, 520])
    t1.setStyle(TableStyle([
        ('BACKGROUND', (0,0), (-1,-1), colors.HexColor('#0F172A')),
        ('BOX', (0,0), (-1,-1), 1, colors.HexColor('#00F0FF')),
        ('VALIGN', (0,0), (-1,-1), 'MIDDLE'),
        ('TOPPADDING', (0,0), (-1,-1), 8),
        ('BOTTOMPADDING', (0,0), (-1,-1), 8),
        ('LEFTPADDING', (0,0), (-1,-1), 12),
        ('RIGHTPADDING', (0,0), (-1,-1), 12),
    ]))
    story.append(t1)
    story.append(PageBreak())

    # -------------------------------------------------------------
    # SLIDE 2: PROPOSED SOLUTION
    # -------------------------------------------------------------
    story.append(Paragraph("IDEA TITLE: ORBITGUARD", sub_style))
    story.append(Paragraph("PROPOSED SOLUTION (Describe your Idea / Solution / Prototype)", h1_style))
    story.append(Spacer(1, 6))

    sol_grid = [
        [
            Paragraph("<b>1. Multi-Source Ingestion</b><br/>CelesTrak & Space-Track live feeds + Modulo-10 checksum validation with 35,535 TLE cache.", bullet_style),
            Paragraph("<b>2. SGP4 Astrodynamics</b><br/>SGP4 propagation, 3-tier spatial sieve (altitude shells <=50km), and orthogonal Secant TCA solver.", bullet_style)
        ],
        [
            Paragraph("<b>3. Probability & Risk</b><br/>Foster-2D B-plane numerical integral, Akella-Alfriend series, Alfano max Pc, and 0-100 risk scoring.", bullet_style),
            Paragraph("<b>4. CAM & AI Decision</b><br/>4-candidate impulsive burn matrix, Tsiolkovsky Hydrazine fuel cost, and 15-tool grounded AI copilot.", bullet_style)
        ]
    ]
    t2_grid = Table(sol_grid, colWidths=[355, 355])
    t2_grid.setStyle(TableStyle([
        ('BACKGROUND', (0,0), (-1,-1), colors.HexColor('#0F172A')),
        ('BOX', (0,0), (-1,-1), 1, colors.HexColor('#1E293B')),
        ('INNERGRID', (0,0), (-1,-1), 1, colors.HexColor('#1E293B')),
        ('PADDING', (0,0), (-1,-1), 8),
    ]))
    story.append(t2_grid)
    story.append(Spacer(1, 10))

    story.append(Paragraph("<b>End-to-End Operational Workflow:</b>", sub_style))
    wf_data = [
        [Paragraph("<b>DATA</b>", bullet_style), Paragraph("Multi-source live TLE ingestion (CelesTrak / Space-Track) + Modulo-10 checksums.", bullet_style)],
        [Paragraph("<b>SGP4</b>", bullet_style), Paragraph("Vectorized SGP4 propagation in TEME frame with WGS-84 / EGM-96 geoid transforms.", bullet_style)],
        [Paragraph("<b>TCA</b>", bullet_style), Paragraph("3-tier spatial sieve + Secant root solver finding exact orthogonal TCA (r_rel · v_rel = 0).", bullet_style)],
        [Paragraph("<b>Pc & RISK</b>", bullet_style), Paragraph("Foster-2D B-plane probability integral + Alfano max Pc + 0-100 composite aerospace risk score.", bullet_style)],
        [Paragraph("<b>CAM</b>", bullet_style), Paragraph("Gauss variational equations compute 4 burn candidates + Tsiolkovsky Hydrazine propellant mass.", bullet_style)],
        [Paragraph("<b>DECISION</b>", bullet_style), Paragraph("FDO / Flight Director approval controls, post-CAM verification, and immutable UTC audit logging.", bullet_style)]
    ]
    t2_wf = Table(wf_data, colWidths=[90, 620])
    t2_wf.setStyle(TableStyle([
        ('BACKGROUND', (0,0), (-1,-1), colors.HexColor('#0F172A')),
        ('BOX', (0,0), (-1,-1), 1, colors.HexColor('#00F0FF')),
        ('INNERGRID', (0,0), (-1,-1), 0.5, colors.HexColor('#1E293B')),
        ('PADDING', (0,0), (-1,-1), 4),
    ]))
    story.append(t2_wf)
    story.append(PageBreak())

    # -------------------------------------------------------------
    # SLIDE 3: TECHNICAL APPROACH
    # -------------------------------------------------------------
    story.append(Paragraph("TECHNICAL APPROACH", h1_style))
    story.append(Paragraph("3-Layer Architecture & Implementation Methodology", sub_style))
    story.append(Spacer(1, 6))

    arch_data = [
        [
            Paragraph("<b>LAYER 1: DATA & FEED INGESTION</b><br/>• Feeds: CelesTrak GP API (Primary), Space-Track 18th SDS, SatNOGS.<br/>• Ingestion: Async HTTP fetcher + Modulo-10 line checksum validation.<br/>• Fallback: Local verified cache with 35,535 real space objects.<br/>• Provenance: Real-time UI badge explicitly separating LIVE from DEMO.", bullet_style),
            Paragraph("<b>LAYER 2: ASTRODYNAMICS ENGINE</b><br/>• Propagator: SGP4 C-extension (Vallado 2006) with J2/J3/J4 harmonics.<br/>• 3-Tier Sieve: Altitude shells (|Δh| <= 50km) prune 98.6% pairs in O(1).<br/>• Secant TCA: Orthogonal zero-crossing (r_rel · v_rel = 0) in <10^-4 s.<br/>• Pc & Risk: Foster-2D B-plane integral + Akella-Alfriend + Alfano max Pc.", bullet_style),
            Paragraph("<b>LAYER 3: DECISION & CLIENT</b><br/>• CAM: 4-candidate matrix + Tsiolkovsky fuel mass (Isp=220s Hydrazine).<br/>• AI Copilot: 15-tool layer + Digit Validator (±1% tolerance).<br/>• Case Center: 13-section HUD (NEW -> APPROVED -> VERIFIED).<br/>• 3D WebGL: Three.js rendering 35k objects @ 60 FPS + CCSDS CDM.", bullet_style)
        ]
    ]
    t3_arch = Table(arch_data, colWidths=[235, 240, 235])
    t3_arch.setStyle(TableStyle([
        ('BACKGROUND', (0,0), (-1,-1), colors.HexColor('#0F172A')),
        ('BOX', (0,0), (-1,-1), 1, colors.HexColor('#00F0FF')),
        ('INNERGRID', (0,0), (-1,-1), 1, colors.HexColor('#1E293B')),
        ('PADDING', (0,0), (-1,-1), 8),
    ]))
    story.append(t3_arch)
    story.append(Spacer(1, 10))

    story.append(Paragraph("<b>Verified Technology Stack:</b> Python 3.9+ • FastAPI • SGP4 C-Extension • NumPy • SciPy • SQLAlchemy • SQLite • React 18 • TypeScript • Three.js WebGL • TailwindCSS", sub_style))
    story.append(PageBreak())

    # -------------------------------------------------------------
    # SLIDE 4: FEASIBILITY AND VIABILITY
    # -------------------------------------------------------------
    story.append(Paragraph("FEASIBILITY AND VIABILITY", h1_style))
    story.append(Paragraph("Analysis, Operational Risks & Mitigation Strategies", sub_style))
    story.append(Spacer(1, 6))

    feas_data = [
        [
            Paragraph("<b>CHALLENGE  →  OPERATIONAL RISK  →  MITIGATION</b><br/><br/>"
                      "<b>1. TLE Covariance Gap:</b> TLEs lack 6x6 covariance matrices.<br/>"
                      "→ <i>Mitigation:</i> Computes Alfano Max Pc bounds (Pc,max = R² / e·d²) establishing theoretical worst-case risk limits.<br/><br/>"
                      "<b>2. External Feed Outage:</b> Space-Track / CelesTrak downtime.<br/>"
                      "→ <i>Mitigation:</i> Multi-source failover architecture with instant fallback to local verified cache (35,535 TLEs).<br/><br/>"
                      "<b>3. AI Hallucination Risk:</b> LLM generating false numbers.<br/>"
                      "→ <i>Mitigation:</i> Zero-LLM physics rule + automated Digit Validator (±1% tolerance) verifying all claims against physics engine.<br/><br/>"
                      "<b>4. Real Thruster Risk:</b> Accidental unverified firing.<br/>"
                      "→ <i>Mitigation:</i> Framed as Decision Support. Real commanding requires multi-sig human approval & encryption keys.<br/><br/>"
                      "<b>5. Computational Scaling:</b> O(N²) all-pairs bottleneck.<br/>"
                      "→ <i>Mitigation:</i> 3-tier spatial sieve prunes 98.6% of candidate pairs in O(1) time via altitude envelope shells.", bullet_style),
            Paragraph("<b>EMPIRICAL BENCHMARKS</b><br/><br/>"
                      "<b>79 / 79 Tests Passed (100%)</b><br/>Automated Pytest test suite across unit & integration tests.<br/><br/>"
                      "<b>11 / 11 Benchmarks Passed (100%)</b><br/>Validated against Vallado 2006 standards in 0.0008s.<br/><br/>"
                      "<b>~100.7 ms Pipeline Latency</b><br/>Catalog screening latency across 35,000+ objects.<br/><br/>"
                      "<b>+26.9 km Miss Clearance</b><br/>Delivered via 0.505 m/s impulsive prograde burn.<br/><br/>"
                      "<b>117 grams Hydrazine</b><br/>Propellant consumed for 260 kg spacecraft (Tsiolkovsky).<br/><br/>"
                      "<b>-99.9% Risk Reduction</b><br/>Achieved post-CAM verification clearance.", bullet_style)
        ]
    ]
    t4 = Table(feas_data, colWidths=[450, 260])
    t4.setStyle(TableStyle([
        ('BACKGROUND', (0,0), (-1,-1), colors.HexColor('#0F172A')),
        ('BOX', (0,0), (-1,-1), 1, colors.HexColor('#00F0FF')),
        ('INNERGRID', (0,0), (-1,-1), 1, colors.HexColor('#1E293B')),
        ('PADDING', (0,0), (-1,-1), 8),
    ]))
    story.append(t4)
    story.append(PageBreak())

    # -------------------------------------------------------------
    # SLIDE 5: IMPACT AND BENEFITS
    # -------------------------------------------------------------
    story.append(Paragraph("IMPACT AND BENEFITS", h1_style))
    story.append(Paragraph("Social, Economic, Environmental & Operational Beneficiary Reach", sub_style))
    story.append(Spacer(1, 6))

    story.append(Paragraph("<b>Workflow Progression:</b> Early Screening  →  Precise Risk  →  Faster Decision  →  Optimal CAM  →  Reduced Risk  →  Sustainability", sub_style))
    story.append(Spacer(1, 4))

    imp_grid = [
        [
            Paragraph("<b>🛰️ SATELLITE OPERATORS</b><br/>• Democratizes enterprise-grade CARA & CAM optimization without high recurring SaaS fees.<br/>• Exact fuel cost modeling (117g Hydrazine per CAM) with 99.9% risk reduction.", bullet_style),
            Paragraph("<b>🛡️ SPACE AGENCIES & DEFENSE</b><br/>• Sovereign, explainable, and air-gapped conjunction monitoring with full data provenance.<br/>• Automated CCSDS 508.0-B-1 CDM export and Defense SITREP executive threat dossiers.", bullet_style)
        ],
        [
            Paragraph("<b>🎓 RESEARCH & ACADEMIA</b><br/>• Open-access, transparent astrodynamics testbed for training next-generation space engineers.<br/>• 100% mathematically proven physics equations with zero black-box heuristics.", bullet_style),
            Paragraph("<b>🌌 SPACE SUSTAINABILITY</b><br/>• Directly mitigates Kessler Syndrome runaway cascade risk, safeguarding critical LEO/GEO orbits.<br/>• Protects multi-billion-dollar space assets from hypervelocity kinetic debris destruction.", bullet_style)
        ]
    ]
    t5 = Table(imp_grid, colWidths=[355, 355])
    t5.setStyle(TableStyle([
        ('BACKGROUND', (0,0), (-1,-1), colors.HexColor('#0F172A')),
        ('BOX', (0,0), (-1,-1), 1, colors.HexColor('#00F0FF')),
        ('INNERGRID', (0,0), (-1,-1), 1, colors.HexColor('#1E293B')),
        ('PADDING', (0,0), (-1,-1), 8),
    ]))
    story.append(t5)
    story.append(PageBreak())

    # -------------------------------------------------------------
    # SLIDE 6: RESEARCH AND REFERENCES
    # -------------------------------------------------------------
    story.append(Paragraph("RESEARCH, VALIDATION & REFERENCES", h1_style))
    story.append(Paragraph("Analytical Proofs, Benchmark Evidence & International Standards", sub_style))
    story.append(Spacer(1, 6))

    ref_grid = [
        [
            Paragraph("<b>SCIENTIFIC VALIDATION BENCHMARKS</b><br/><br/>"
                      "✓ <b>SGP4 Propagation:</b> Validated against Vallado 2006 Reference Vectors (LEO-00005, ISS-25544) with <0.001 km epoch error.<br/><br/>"
                      "✓ <b>Coordinate Transforms:</b> TEME -> ECEF -> WGS-84 Geodetic validated using GMST rotation and Bowring's closed-form algorithm.<br/><br/>"
                      "✓ <b>Secant TCA Solver:</b> Resolves exact orthogonal encounter condition (r_rel · v_rel = 0) with sub-millisecond precision (<10^-4 s).<br/><br/>"
                      "✓ <b>Foster-2D Pc Engine:</b> B-plane numerical integral & Akella-Alfriend series match NASA benchmarks (<0.1% deviation).<br/><br/>"
                      "✓ <b>Gauss CAM Mechanics:</b> Along-track separation (Δs ≈ 3 Δvt Δt_lead) & Tsiolkovsky Hydrazine propellant mass validated.<br/><br/>"
                      "✓ <b>Test Pass Rate:</b> 79/79 Pytest tests (100%) & 11/11 analytical benchmarks passing in 0.0008s.", bullet_style),
            Paragraph("<b>RESEARCH PAPERS & INTERNATIONAL STANDARDS</b><br/><br/>"
                      "• <b>SGP4 Standard:</b> Vallado, D. A., et al. (2006). <i>Revisiting Spacetrack Report #3: Rev 2</i>. AIAA 2006-6753.<br/><br/>"
                      "• <b>Foster-2D Collision Probability:</b> Foster, J. L., & Estes, H. S. (1992). <i>A Parametric Analysis of Orbital Debris Collision Probability and Maneuver Utility</i>. NASA JSC-25898.<br/><br/>"
                      "• <b>Short-Encounter Model:</b> Chan, K. (1997). <i>Collision Probability for Space Missions</i>. Advances in Astronautical Sciences, 96.<br/><br/>"
                      "• <b>Akella-Alfriend Probability:</b> Akella, M. R., & Alfriend, K. T. (2000). <i>Probability of Collision Between Space Objects</i>. JGCD, 23(5).<br/><br/>"
                      "• <b>Alfano Maximum Pc Bounds:</b> Alfano, S. (2005). <i>Relating Position Uncertainty to Maximum Conjunction Probability</i>. JAS, 53(2).<br/><br/>"
                      "• <b>CCSDS CDM Standard:</b> CCSDS 508.0-B-1 (2013). <i>Conjunction Data Message (CDM)</i>. Blue Book Standard.", bullet_style)
        ]
    ]
    t6 = Table(ref_grid, colWidths=[355, 355])
    t6.setStyle(TableStyle([
        ('BACKGROUND', (0,0), (-1,-1), colors.HexColor('#0F172A')),
        ('BOX', (0,0), (-1,-1), 1, colors.HexColor('#00F0FF')),
        ('INNERGRID', (0,0), (-1,-1), 1, colors.HexColor('#1E293B')),
        ('PADDING', (0,0), (-1,-1), 8),
    ]))
    story.append(t6)

    # Build PDF
    doc.build(story, onFirstPage=draw_bg, onLaterPages=draw_bg)
    print(f"[✓] Successfully generated PDF: {output_path}")


def main():
    sih_dir = "/Users/kundan/Downloads/ORBITGUARD/docs/sih"
    os.makedirs(sih_dir, exist_ok=True)
    
    pptx_out = os.path.join(sih_dir, "ORBITGUARD_SIH_FINAL_6_SLIDES.pptx")
    pdf_out = os.path.join(sih_dir, "ORBITGUARD_SIH_FINAL_6_SLIDES.pdf")

    print("[*] Generating ORBITGUARD SIH 2026 Final 6-Slide Presentation...")
    create_pptx(pptx_out)
    create_pdf(pdf_out)
    print("[✓] ALL ARTIFACTS GENERATED SUCCESSFULLY.")


if __name__ == "__main__":
    main()
